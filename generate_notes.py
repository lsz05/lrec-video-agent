"""
generate_notes.py

Generates English speaker notes for each slide in an English .pptx using Claude,
based on the slide's text content and an optional paper PDF for context.
Useful when slides have no existing notes, or to regenerate notes from scratch.

Usage:
    python3 generate_notes.py <slides_en.pptx> [--paper <paper.pdf>]
                              [--target-wpm <wpm>] [--target-min <minutes>]
                              [--overwrite]

    <slides_en.pptx>    English .pptx file (output of translate_slides.py)
    --paper             Paper PDF for context — improves terminology and accuracy
    --target-wpm        Speaking speed in words per minute (default: 130)
    --target-min        Target total presentation duration in minutes (default: 15)
    --overwrite         Overwrite existing notes (default: skip slides that already have notes)

Output:
    Updates <slides_en.pptx> in-place with generated speaker notes.
    Backup saved as <slides_en>_before_notes.pptx.

Requirements:
    pip install anthropic python-pptx PyMuPDF
    export ANTHROPIC_API_KEY=sk-ant-...
"""

import argparse
import json
import os
import shutil
import sys
import time
from pathlib import Path

import anthropic
from pptx import Presentation


# ---------------------------------------------------------------------------
# PDF context extraction
# ---------------------------------------------------------------------------

def extract_pdf_context(pdf_path: Path, max_chars: int = 4000) -> str:
    """Extract abstract/introduction from a PDF for use as context."""
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) >= max_chars:
                break
        doc.close()
        return text[:max_chars]
    except ImportError:
        print("  [warn] PyMuPDF not installed — skipping paper context.", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"  [warn] Could not read PDF: {e}", file=sys.stderr)
        return ""


# ---------------------------------------------------------------------------
# Slide text extraction
# ---------------------------------------------------------------------------

def extract_slide_text(slide) -> str:
    """Extract all visible text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
        elif shape.has_table:
            for row in shape.table.rows:
                row_text = "  |  ".join(cell.text_frame.text.strip() for cell in row.cells)
                if row_text.strip("  |  "):
                    texts.append(row_text)
    return "\n".join(texts)


# ---------------------------------------------------------------------------
# Note generation via Claude
# ---------------------------------------------------------------------------

MAX_RETRIES = 3
RETRY_BACKOFF = [2, 5, 15]


def generate_note(
    client: anthropic.Anthropic,
    slide_number: int,
    total_slides: int,
    slide_text: str,
    target_words: int,
    paper_context: str,
) -> str:
    """Generate a speaker note for a single slide using Claude."""

    context_block = ""
    if paper_context:
        context_block = f"""
Paper context (use for accurate terminology and background):
<paper_context>
{paper_context}
</paper_context>
"""

    prompt = f"""You are writing speaker notes for an academic conference presentation.
{context_block}
Generate a speaker note for slide {slide_number} of {total_slides}.
Target length: approximately {target_words} words.

Rules:
- Write in first person (e.g. "In this slide, I will show...")
- Natural spoken English — not bullet points, no headings
- Match the slide content accurately; do not invent facts
- Provide context and explanation that goes beyond what is written on the slide
- End with a smooth transition to the next slide if appropriate (except the last slide)
- Return ONLY the speaker note text, no extra commentary

Slide content:
{slide_text}"""

    last_error = None
    for attempt in range(MAX_RETRIES):
        try:
            response = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=1024,
                messages=[{"role": "user", "content": prompt}],
            )
            return response.content[0].text.strip()
        except anthropic.RateLimitError as e:
            last_error = str(e)
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Rate limit. Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)
        except anthropic.APIError as e:
            last_error = str(e)
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] API error: {e}. Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

    raise RuntimeError(f"All retries failed for slide {slide_number}: {last_error}")


# ---------------------------------------------------------------------------
# Notes writing
# ---------------------------------------------------------------------------

def set_notes(slide, text: str) -> None:
    """Write text into slide notes."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    lines = text.split("\n")
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


# ---------------------------------------------------------------------------
# Checkpoint helpers
# ---------------------------------------------------------------------------

def checkpoint_path(pptx_path: Path) -> Path:
    return pptx_path.parent / (pptx_path.stem + "_notes_gen_checkpoint.json")


def load_checkpoint(pptx_path: Path) -> set[int]:
    cp = checkpoint_path(pptx_path)
    if cp.exists():
        try:
            done = set(json.loads(cp.read_text()))
            print(f"  Resuming from checkpoint ({len(done)} slides already generated)")
            return done
        except Exception:
            pass
    return set()


def save_checkpoint(pptx_path: Path, done: set[int]) -> None:
    checkpoint_path(pptx_path).write_text(json.dumps(sorted(done)))


def remove_checkpoint(pptx_path: Path) -> None:
    cp = checkpoint_path(pptx_path)
    if cp.exists():
        cp.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Generate English speaker notes for slides using Claude.")
    parser.add_argument("input", help="Path to the English .pptx")
    parser.add_argument("--paper", help="Paper PDF for context (optional)")
    parser.add_argument("--target-wpm", type=int, default=130,
                        help="Speaking speed in words per minute (default: 130)")
    parser.add_argument("--target-min", type=float, default=15.0,
                        help="Target total duration in minutes (default: 15)")
    parser.add_argument("--overwrite", action="store_true",
                        help="Overwrite existing notes (default: skip slides that already have notes)")
    args = parser.parse_args()

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY is not set.", file=sys.stderr)
        sys.exit(1)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    # Backup
    backup_path = input_path.with_stem(input_path.stem + "_before_notes")
    if not backup_path.exists():
        shutil.copy2(input_path, backup_path)
        print(f"Backup saved: {backup_path}")

    # Paper context
    paper_context = ""
    if args.paper:
        print(f"Extracting context from {args.paper}...")
        paper_context = extract_pdf_context(Path(args.paper))

    client = anthropic.Anthropic()
    prs = Presentation(str(input_path))
    total = len(prs.slides)

    # Calculate per-slide target word count
    # Distribute evenly across slides that will get notes
    total_target_words = int(args.target_min * args.target_wpm)
    slides_needing_notes = []
    for i, slide in enumerate(prs.slides):
        text = extract_slide_text(slide)
        has_notes = (slide.has_notes_slide and
                     slide.notes_slide.notes_text_frame.text.strip())
        if text and (args.overwrite or not has_notes):
            slides_needing_notes.append(i)

    if not slides_needing_notes:
        print("All slides already have notes. Use --overwrite to regenerate.")
        sys.exit(0)

    per_slide_words = max(10, total_target_words // len(slides_needing_notes))
    print(f"Generating notes for {len(slides_needing_notes)}/{total} slides "
          f"(~{per_slide_words} words each, target {args.target_min} min total)\n")

    done_slides = load_checkpoint(input_path)

    try:
        for i, slide in enumerate(prs.slides, start=1):
            slide_idx = i - 1
            if slide_idx not in slides_needing_notes:
                text = extract_slide_text(slide)
                reason = "(no text)" if not text else "(already has notes)"
                print(f"  Slide {i:02d}/{total}  {reason} — skipped")
                continue

            if i in done_slides:
                print(f"  Slide {i:02d}/{total}  (from checkpoint)")
                continue

            slide_text = extract_slide_text(slide)
            print(f"  Slide {i:02d}/{total}  ({len(slide_text.split())} words on slide)...",
                  end=" ", flush=True)

            note = generate_note(
                client, i, total, slide_text, per_slide_words, paper_context)
            set_notes(slide, note)

            done_slides.add(i)
            save_checkpoint(input_path, done_slides)
            prs.save(str(input_path))
            print(f"done ({len(note.split())} words)")

    except KeyboardInterrupt:
        print("\n\nInterrupted. Progress saved — re-run to resume.", file=sys.stderr)
        sys.exit(1)

    remove_checkpoint(input_path)

    # Final word count
    total_words = sum(
        len(s.notes_slide.notes_text_frame.text.strip().split())
        for s in prs.slides
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
    )
    print(f"\nDone. Total notes: {total_words} words "
          f"(~{total_words / args.target_wpm:.1f} min at {args.target_wpm} wpm)")
    print(f"Saved: {input_path}")


if __name__ == "__main__":
    main()
