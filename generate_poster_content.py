"""
generate_poster_content.py

Reads a paper PDF and the translated slides (.pptx), then uses Claude to
produce a structured JSON file containing all content needed to render an
academic poster.

Usage:
    python generate_poster_content.py <paper.pdf> <slides_en.pptx> [--out <output.json>]

    <paper.pdf>       Paper PDF
    <slides_en.pptx>  English slides produced by translate_slides.py / translate_notes.py
    --out             Output JSON path (default: <paper>_poster_content.json)

Requirements:
    pip install anthropic python-pptx PyMuPDF
    export ANTHROPIC_API_KEY=sk-ant-...
"""

import argparse
import json
import os
import sys
from pathlib import Path

import anthropic
from pptx import Presentation


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------

def extract_pdf_text(pdf_path: str) -> str:
    try:
        import fitz
    except ImportError:
        print("Error: PyMuPDF not installed. Run: pip install PyMuPDF", file=sys.stderr)
        sys.exit(1)
    doc = fitz.open(pdf_path)
    return "\n".join(page.get_text() for page in doc)


def extract_slides_text(pptx_path: str) -> str:
    """Extract slide text and speaker notes as a readable string."""
    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        notes = (
            slide.notes_slide.notes_text_frame.text.strip()
            if slide.has_notes_slide else ""
        )
        if texts or notes:
            lines.append(f"=== Slide {i} ===")
            lines.extend(texts)
            if notes:
                lines.append(f"[Speaker notes: {notes}]")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Poster content generation via Claude
# ---------------------------------------------------------------------------

SYSTEM_PROMPT = """\
You are an expert academic poster designer specializing in NLP / AI research.
Given a research paper and its presentation slides, extract and condense the
content into a richly detailed academic poster for LREC 2026.

A good academic poster is DENSE with information:
- Uses short, punchy bullet points (10–25 words each). No full paragraphs.
- Every results/analysis bullet MUST include specific numbers (scores, counts,
  percentages, speedups, correlation values, model names, dataset names, etc.).
- Do not use vague claims like "significantly better" — always give the number.
- Has a single memorable takeaway message.
- Is organized into 5–6 sections that flow logically.
- Includes standout key numbers, key tables, and key figures from the paper.

Return ONLY a JSON object with this exact schema (no explanation, no markdown fences):
{
  "title": "<full paper title>",
  "authors": "<Author1, Author2, ...>",
  "affiliations": "<Affil1; Affil2>",
  "conference": "<conference name and year>",
  "key_takeaway": "<one sentence — the single most important contribution or result>",
  "key_numbers": [
    {"stat": "<number or short phrase>", "label": "<what it measures>"},
    ...
  ],
  "key_tables": [
    {
      "caption": "<short table caption>",
      "headers": ["<col1>", "<col2>", ...],
      "rows": [["<cell>", ...], ...]
    }
  ],
  "key_figures": [
    {
      "slide": <1-based slide number in the slides file>,
      "caption": "<short descriptive caption for the figure>"
    }
  ],
  "sections": [
    {
      "heading": "<section title>",
      "bullets": ["<bullet 1>", "<bullet 2>", ...]
    },
    ...
  ]
}

Guidelines:
- "key_numbers": 6–8 entries with the most impressive concrete statistics.
- "key_tables": 1–2 tables. Keep narrow (≤5 columns, ≤12 rows). Pick the most
  informative table (e.g. top-10 leaderboard, corpus reduction stats).
  Abbreviate long model names to fit (e.g. "sarashina-v2-1b").
- "key_figures": 1–2 figures. Identify the slide number (1-based) from the
  slides file that contains the most informative figure (scatter plots, bar
  charts, pipeline diagrams, result tables as images). Skip title, agenda,
  and text-only slides. Prefer slides that are primarily visual.
- "sections": 5–6 sections, each with 6–8 specific, number-rich bullets.
- Results bullets must cite actual numbers from the paper's tables and figures.
- "key_takeaway": one sentence suitable for a large banner on the poster.
"""


def generate_content(
    client: anthropic.Anthropic,
    paper_text: str,
    slides_text: str,
) -> dict:
    user_content = (
        "PAPER:\n" + paper_text[:12000]
        + "\n\n---\n\nSLIDES:\n" + slides_text[:6000]
    )

    last_error = None
    for attempt in range(3):
        try:
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=4096,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": user_content}],
            )
            raw = message.content[0].text.strip()
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            return json.loads(raw.strip())

        except json.JSONDecodeError as e:
            last_error = f"JSON parse error: {e}"
            print(f"  [warn] Attempt {attempt + 1}/3 — {last_error}", file=sys.stderr)

        except anthropic.APIStatusError as e:
            last_error = f"API error {e.status_code}"
            if e.status_code < 500:
                break
            print(f"  [warn] Attempt {attempt + 1}/3 — {last_error}", file=sys.stderr)

    print(f"Error: failed to generate poster content ({last_error})", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Generate structured poster content JSON from a paper and slides.")
    parser.add_argument("paper", help="Path to paper PDF")
    parser.add_argument("slides", help="Path to English slides .pptx")
    parser.add_argument("--out", help="Output JSON path")
    args = parser.parse_args()

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY is not set.", file=sys.stderr)
        sys.exit(1)

    paper_path = Path(args.paper)
    slides_path = Path(args.slides)
    for p in [paper_path, slides_path]:
        if not p.exists():
            print(f"Error: file not found: {p}", file=sys.stderr)
            sys.exit(1)

    out_path = Path(args.out) if args.out else paper_path.with_stem(paper_path.stem + "_poster_content").with_suffix(".json")

    client = anthropic.Anthropic()

    print(f"Extracting paper text from {paper_path.name}...")
    paper_text = extract_pdf_text(str(paper_path))

    print(f"Extracting slides text from {slides_path.name}...")
    slides_text = extract_slides_text(str(slides_path))

    print("Generating poster content via Claude...")
    content = generate_content(client, paper_text, slides_text)

    out_path.write_text(json.dumps(content, ensure_ascii=False, indent=2))
    print(f"\nSaved: {out_path}")
    print(f"  Sections : {len(content.get('sections', []))}")
    for s in content.get("sections", []):
        print(f"    - {s['heading']} ({len(s['bullets'])} bullets)")
    print(f"  Tables   : {len(content.get('key_tables', []))}")
    for t in content.get("key_tables", []):
        print(f"    - {t['caption']} ({len(t['headers'])} cols × {len(t['rows'])} rows)")
    print(f"  Figures  : {len(content.get('key_figures', []))}")
    for f in content.get("key_figures", []):
        print(f"    - Slide {f['slide']}: {f['caption'][:60]}")
    print(f"  Takeaway : {content.get('key_takeaway', '')}")


if __name__ == "__main__":
    main()
