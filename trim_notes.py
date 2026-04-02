"""
trim_notes.py

Trims speaker notes in an English .pptx to fit a target presentation duration,
using Claude to shorten each note while preserving the key message.

Usage:
    python3 trim_notes.py <slides_en.pptx> [--target-min <minutes>] [--wpm <wpm>]

    --target-min    Target presentation duration in minutes (default: 15)
    --wpm           Assumed speaking speed in words per minute (default: 130)

Output:
    Overwrites the input .pptx with trimmed notes (original backed up as <stem>_backup.pptx)

Requirements:
    pip install anthropic python-pptx
    export ANTHROPIC_API_KEY=...
"""

import argparse
import json
import os
import shutil
import sys
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.util import Pt
from lxml import etree


def extract_notes(prs: Presentation) -> list[tuple[int, str]]:
    """Return list of (slide_index, note_text) for slides with non-empty notes."""
    result = []
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                result.append((i, text))
    return result


def set_notes(slide, text: str) -> None:
    """Write text into slide notes, preserving existing formatting structure."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing paragraphs except the first
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)
    # Set first paragraph
    lines = text.split("\n")
    tf.paragraphs[0].text = lines[0]
    # Add additional lines as new paragraphs
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def trim_single_note(
    client: anthropic.Anthropic,
    slide_num: int,
    text: str,
    target_words: int,
) -> str:
    """Trim a single slide's note to approximately target_words words."""
    current_words = len(text.split())
    if current_words <= target_words:
        return text

    prompt = f"""Trim this speaker note from {current_words} to approximately {target_words} words.
Preserve the speaker's first-person voice and all key technical points.
Remove filler phrases, redundant transitions, and over-explanation only.
Return ONLY the trimmed text with no extra commentary.

{text}"""

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.content[0].text.strip()


def trim_notes_with_claude(
    client: anthropic.Anthropic,
    notes: list[tuple[int, str]],
    target_words: int,
) -> dict[int, str]:
    """
    Trim notes proportionally to collectively fit within target_words.
    Returns a dict mapping slide_index -> trimmed_text.
    """
    current_words = sum(len(t.split()) for _, t in notes)
    reduction_ratio = target_words / current_words

    print(f"Trimming {current_words} → {target_words} words ({round((1-reduction_ratio)*100)}% reduction)...")

    result = {}
    for slide_idx, text in notes:
        slide_num = slide_idx + 1
        words = len(text.split())
        slide_target = max(5, round(words * reduction_ratio))
        print(f"  Slide {slide_num:02d}  ({words} → ~{slide_target} words)...", end=" ", flush=True)
        trimmed = trim_single_note(client, slide_num, text, slide_target)
        result[slide_idx] = trimmed
        print(f"done ({len(trimmed.split())} words)")

    return result


def count_words(prs: Presentation) -> int:
    total = 0
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                total += len(text.split())
    return total


def main():
    parser = argparse.ArgumentParser(description="Trim speaker notes to fit a target duration.")
    parser.add_argument("input", help="Path to the English .pptx")
    parser.add_argument("--target-min", type=float, default=15.0,
                        help="Target duration in minutes (default: 15)")
    parser.add_argument("--wpm", type=int, default=130,
                        help="Speaking speed in words per minute (default: 130)")
    args = parser.parse_args()

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY is not set.", file=sys.stderr)
        sys.exit(1)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    target_words = int(args.target_min * args.wpm)
    print(f"Target: {args.target_min} min at {args.wpm} wpm = {target_words} words")

    # Backup original
    backup_path = input_path.with_stem(input_path.stem + "_backup")
    shutil.copy2(input_path, backup_path)
    print(f"Backup saved: {backup_path}")

    prs = Presentation(str(input_path))
    notes = extract_notes(prs)
    current_words = sum(len(t.split()) for _, t in notes)
    print(f"Current: {current_words} words across {len(notes)} slides\n")

    if current_words <= target_words:
        print(f"Already within target ({current_words} ≤ {target_words} words). Nothing to do.")
        sys.exit(0)

    client = anthropic.Anthropic()
    trimmed = trim_notes_with_claude(client, notes, target_words)

    # Apply trimmed notes back to PPTX
    for slide_idx, text in trimmed.items():
        set_notes(prs.slides[slide_idx], text)

    prs.save(str(input_path))

    # Report result
    prs2 = Presentation(str(input_path))
    new_words = count_words(prs2)
    print(f"\nDone. {current_words} → {new_words} words "
          f"(~{new_words / args.wpm:.1f} min at {args.wpm} wpm)")
    print(f"Saved: {input_path}")


if __name__ == "__main__":
    main()
