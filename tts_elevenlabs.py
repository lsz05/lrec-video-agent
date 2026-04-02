"""
tts_elevenlabs.py

Reads English speaker notes from a .pptx and synthesizes each slide's notes
into an audio file using ElevenLabs TTS — supports custom / cloned voices.

Usage:
    python tts_elevenlabs.py <slides_en.pptx> --voice-id <voice_id> [options]

    <slides_en.pptx>    English .pptx produced by translate_notes.py
    --voice-id          ElevenLabs voice ID (required)
                        Find it at https://elevenlabs.io/app/voice-lab
                        After cloning your voice, copy the Voice ID from the card.
    --out-dir           Directory to write audio files into (default: <pptx_stem>_audio/)
    --model             ElevenLabs model (default: eleven_multilingual_v2)
                        Choices: eleven_multilingual_v2, eleven_turbo_v2_5, eleven_flash_v2_5
    --stability         Voice stability 0.0–1.0 (default: 0.5, higher = more consistent)
    --similarity        Similarity boost 0.0–1.0 (default: 0.8, higher = closer to original)
    --style             Style exaggeration 0.0–1.0 (default: 0.0)
    --speaker-boost     Enable speaker boost for more expressive output (flag, default off)

Output:
    One .mp3 per slide: slide_01.mp3, slide_02.mp3, ...
    Slides with no notes produce no file.
    A manifest.json is written listing each slide's audio file and duration.

How to clone your voice with ElevenLabs:
    1. Sign up at https://elevenlabs.io (free tier: 1 instant voice clone)
    2. Go to Voice Lab → Add Generative or Cloned Voice → Instant Voice Clone
    3. Upload 1–5 minutes of clear recording of your voice (WAV or MP3, no background noise)
    4. Give it a name and save
    5. Copy the Voice ID from the voice card (Settings icon → Voice ID)
    6. Set your API key: export ELEVENLABS_API_KEY=sk-...
    7. Run: python tts_elevenlabs.py 1091/1091_Slides_en.pptx --voice-id <your_voice_id>

Requirements:
    pip install elevenlabs python-pptx mutagen
    export ELEVENLABS_API_KEY=sk-...
"""

import argparse
import json
import os
import sys
import time
from pathlib import Path

from pptx import Presentation


# ---------------------------------------------------------------------------
# Notes extraction
# ---------------------------------------------------------------------------

def extract_notes(prs: Presentation) -> list[str]:
    """Return a list of speaker note strings, one per slide (empty string if none)."""
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            text = ""
        notes.append(text)
    return notes


# ---------------------------------------------------------------------------
# Pronunciation substitutions
# ---------------------------------------------------------------------------

# Applied before sending text to ElevenLabs so terms are spoken correctly.
# Longer/more-specific terms must come before shorter ones to avoid partial matches.
DEFAULT_PRONUNCIATIONS = {
    "SentenceBERT": "sentence bert",
    "SimCSE": "sim C-S-E",
    "JMTEB-lite": "J-M-teb-lite",
    "JMTEB": "J-M-teb",
    "MMTEB": "M-M-teb",
    "MTEB": "M-teb",
    "STS": "S-T-S",
    "BERT": "bert",
    "SB Intuitions": "S-B intuitions",
    "LLM": "elelem",
    "Mr.TyDi": "mister tidy",
    "MIRACL": "miracle",
}


def apply_pronunciations(text: str, pronunciations: dict) -> str:
    """Replace technical terms with their phonetic equivalents."""
    for term, phonetic in pronunciations.items():
        text = text.replace(term, phonetic)
    return text


# ---------------------------------------------------------------------------
# Audio duration helper
# ---------------------------------------------------------------------------

def get_mp3_duration(path: Path) -> float:
    """Return duration in seconds of an MP3 file using mutagen."""
    try:
        from mutagen.mp3 import MP3
        return MP3(str(path)).info.length
    except Exception:
        return 0.0


# ---------------------------------------------------------------------------
# TTS via ElevenLabs (with retry)
# ---------------------------------------------------------------------------

MAX_RETRIES = 3
RETRY_BACKOFF = [2, 5, 15]


def synthesize(
    client,
    text: str,
    out_path: Path,
    voice_id: str,
    model: str,
    voice_settings: dict,
    slide_number: int,
) -> bool:
    """
    Synthesize text to speech via ElevenLabs and save as MP3 at out_path.
    Returns True on success, False on total failure.
    """
    from elevenlabs import VoiceSettings
    from elevenlabs.core import ApiError

    settings = VoiceSettings(
        stability=voice_settings["stability"],
        similarity_boost=voice_settings["similarity"],
        style=voice_settings["style"],
        use_speaker_boost=voice_settings["speaker_boost"],
        speed=voice_settings["speed"],
    )

    last_error = None
    for attempt in range(MAX_RETRIES):
        try:
            audio_stream = client.text_to_speech.convert(
                voice_id=voice_id,
                text=text,
                model_id=model,
                voice_settings=settings,
                output_format="mp3_44100_128",
            )
            with open(out_path, "wb") as f:
                for chunk in audio_stream:
                    f.write(chunk)
            return True

        except ApiError as e:
            status = getattr(e, "status_code", None)
            last_error = f"API error {status}: {e}"
            if status is not None and status < 500 and status != 429:
                break  # 4xx (not rate-limit) won't fix itself
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except Exception as e:
            last_error = str(e)
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

    print(f"\n    [error] Slide {slide_number}: all retries failed ({last_error}). Skipping.",
          file=sys.stderr)
    return False


# ---------------------------------------------------------------------------
# Checkpoint helpers
# ---------------------------------------------------------------------------

def checkpoint_path(out_dir: Path) -> Path:
    return out_dir / "tts_checkpoint.json"


def load_checkpoint(out_dir: Path) -> set[int]:
    cp = checkpoint_path(out_dir)
    if cp.exists():
        try:
            done = set(json.loads(cp.read_text()))
            print(f"  Resuming from checkpoint ({len(done)} slides already synthesized)")
            return done
        except Exception as e:
            print(f"  [warn] Could not read checkpoint: {e} — starting fresh", file=sys.stderr)
    return set()


def save_checkpoint(out_dir: Path, done: set[int]) -> None:
    checkpoint_path(out_dir).write_text(json.dumps(sorted(done)))


def remove_checkpoint(out_dir: Path) -> None:
    cp = checkpoint_path(out_dir)
    if cp.exists():
        cp.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Synthesize English speaker notes to MP3 using ElevenLabs (supports voice cloning).")
    parser.add_argument("input", help="Path to the English .pptx (with speaker notes)")
    parser.add_argument("--voice-id", required=True,
                        help="ElevenLabs Voice ID (from voice-lab or pre-made voices)")
    parser.add_argument("--out-dir", help="Output directory for audio files")
    parser.add_argument("--model", default="eleven_multilingual_v2",
                        choices=["eleven_multilingual_v2", "eleven_turbo_v2_5", "eleven_flash_v2_5"],
                        help="ElevenLabs model (default: eleven_multilingual_v2)")
    parser.add_argument("--stability", type=float, default=0.5,
                        help="Voice stability 0.0–1.0 (default: 0.5)")
    parser.add_argument("--similarity", type=float, default=0.8,
                        help="Similarity boost 0.0–1.0 (default: 0.8)")
    parser.add_argument("--style", type=float, default=0.0,
                        help="Style exaggeration 0.0–1.0 (default: 0.0)")
    parser.add_argument("--speaker-boost", action="store_true",
                        help="Enable speaker boost for more expressive output")
    parser.add_argument("--speed", type=float, default=1.0,
                        help="Speaking speed multiplier 0.7–1.2 (default: 1.0)")
    parser.add_argument("--max-slides", type=int, default=None,
                        help="Only synthesize the first N slides (useful for testing)")
    args = parser.parse_args()

    api_key = os.environ.get("ELEVENLABS_API_KEY")
    if not api_key:
        print("Error: ELEVENLABS_API_KEY is not set.\n"
              "  export ELEVENLABS_API_KEY=sk-...", file=sys.stderr)
        sys.exit(1)

    try:
        from elevenlabs import ElevenLabs
    except ImportError:
        print("Error: elevenlabs package not installed.\n"
              "  pip install elevenlabs", file=sys.stderr)
        sys.exit(1)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    out_dir = Path(args.out_dir) if args.out_dir else input_path.parent / (input_path.stem + "_audio")
    out_dir.mkdir(parents=True, exist_ok=True)

    client = ElevenLabs(api_key=api_key)

    voice_settings = {
        "stability": args.stability,
        "similarity": args.similarity,
        "style": args.style,
        "speaker_boost": args.speaker_boost,
        "speed": args.speed,
    }

    print(f"Loading {input_path.name}...")
    prs = Presentation(str(input_path))
    notes = extract_notes(prs)
    total = len(notes)
    notes_with_text = sum(1 for n in notes if n.strip())
    print(f"Found {notes_with_text}/{total} slides with notes.")
    print(f"Voice ID: {args.voice_id}  Model: {args.model}  Output: {out_dir}/")
    print(f"Settings: stability={args.stability}  similarity={args.similarity}  "
          f"style={args.style}  speaker_boost={args.speaker_boost}  speed={args.speed}\n")

    done_slides = load_checkpoint(out_dir)
    manifest = []

    try:
        for i, note in enumerate(notes, start=1):
            if args.max_slides is not None and i > args.max_slides:
                break
            entry = {"slide": i, "audio_file": None, "duration_sec": 0.0, "text": note}

            if not note.strip():
                print(f"  Slide {i:02d}/{total}  (no notes — skipped)")
                manifest.append(entry)
                continue

            out_file = out_dir / f"slide_{i:02d}.mp3"

            if i in done_slides and out_file.exists():
                print(f"  Slide {i:02d}/{total}  (from checkpoint)")
                entry["audio_file"] = str(out_file)
                entry["duration_sec"] = get_mp3_duration(out_file)
                manifest.append(entry)
                continue

            spoken = apply_pronunciations(note, DEFAULT_PRONUNCIATIONS)
            print(f"  Slide {i:02d}/{total}  ({len(note)} chars)...", end=" ", flush=True)
            success = synthesize(client, spoken, out_file, args.voice_id, args.model, voice_settings, i)

            if success:
                duration = get_mp3_duration(out_file)
                entry["audio_file"] = str(out_file)
                entry["duration_sec"] = duration
                done_slides.add(i)
                save_checkpoint(out_dir, done_slides)
                print(f"done  ({duration:.1f}s)")
            else:
                print("FAILED")

            manifest.append(entry)

    except KeyboardInterrupt:
        print("\n\nInterrupted. Progress saved — re-run to resume.", file=sys.stderr)
        _write_manifest(out_dir, manifest)
        sys.exit(1)

    _write_manifest(out_dir, manifest)
    remove_checkpoint(out_dir)

    total_duration = sum(e["duration_sec"] for e in manifest)
    print(f"\nTotal audio duration: {total_duration / 60:.1f} min  ({total_duration:.0f}s)")
    if total_duration > 0:
        target = 15 * 60
        diff = total_duration - target
        if abs(diff) > 30:
            direction = "over" if diff > 0 else "under"
            print(f"  [info] {abs(diff):.0f}s {direction} the 15-minute target — "
                  "consider trimming or expanding the notes.")
    print(f"Manifest: {out_dir / 'manifest.json'}")


def _write_manifest(out_dir: Path, manifest: list[dict]) -> None:
    (out_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
