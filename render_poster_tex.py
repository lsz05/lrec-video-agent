"""
render_poster_tex.py

Renders a poster_content.json into a LaTeX/PDF academic poster (A0 portrait).
Uses beamerposter + tcolorbox for a clean, professional layout.

Layout:
  ┌──────────────────────────────────────────┐
  │  HEADER  title / authors / affiliations  │
  ├──────────────────────────────────────────┤
  │  KEY TAKEAWAY (accent banner)            │
  ├──────────────────────────────────────────┤
  │  KEY NUMBERS STRIP                       │
  ├───────────────┬───────────────┬──────────┤
  │  LEFT         │  MIDDLE       │  RIGHT   │
  │  sections 1–3 │  sections 4–6 │  figures │
  │               │               │  tables  │
  └───────────────┴───────────────┴──────────┘

Usage:
    python render_poster_tex.py <poster_content.json>
        [--slides <slides_en.pptx>]
        [--out <poster.tex>]
        [--no-compile]

    --slides      Slides .pptx for figure extraction
    --out         Output .tex path (default: <content stem>_poster.tex)
    --no-compile  Generate .tex only, skip pdflatex

Requirements:
    pip install python-pptx spire.presentation
    pdflatex (MacTeX / TeX Live)
"""

import argparse
import json
import math
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


# ---------------------------------------------------------------------------
# LaTeX character escaping
# ---------------------------------------------------------------------------

_ESCAPE_MAP = str.maketrans({
    "&":  r"\&",
    "%":  r"\%",
    "$":  r"\$",
    "#":  r"\#",
    "_":  r"\_",
    "{":  r"\{",
    "}":  r"\}",
    "~":  r"\textasciitilde{}",
    "^":  r"\^{}",
    "\\": r"\textbackslash{}",
    "▸":  r"$\blacktriangleright$",
    "✦":  r"$\star$",
    "🦴": "",
    "📝": "",
    "📚": "",
    "🤝": "",
    "⏳": "",
    "🫠": "",
})

def esc(text: str) -> str:
    return str(text).translate(_ESCAPE_MAP)


# ---------------------------------------------------------------------------
# Figure extraction from slides (reused from render_poster.py)
# ---------------------------------------------------------------------------

def extract_slide_image(pptx_path: str, slide_number: int, out_path: Path) -> bool:
    """Extract the largest embedded picture from a slide, save to out_path. Returns success."""
    # First try embedded bitmaps via python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        if 1 <= slide_number <= len(prs.slides):
            slide = prs.slides[slide_number - 1]
            pictures = [s for s in slide.shapes if s.shape_type == 13]
            pictures.sort(key=lambda s: s.width * s.height, reverse=True)
            for pic in pictures:
                try:
                    blob = pic.image.blob
                    if blob:
                        out_path.write_bytes(blob)
                        return True
                except (ValueError, AttributeError):
                    continue
    except Exception:
        pass

    # Fall back to Spire for full slide render
    try:
        from spire.presentation import Presentation as SpirePrs
        prs2 = SpirePrs()
        prs2.LoadFromFile(pptx_path)
        slide = prs2.Slides[slide_number - 1]
        img = slide.SaveAsImage()
        img.Save(str(out_path))
        img.Dispose()
        prs2.Dispose()
        return out_path.exists()
    except Exception as e:
        print(f"  [warn] Could not extract slide {slide_number}: {e}", file=sys.stderr)
        return False


# ---------------------------------------------------------------------------
# LaTeX template
# ---------------------------------------------------------------------------

LATEX_PREAMBLE = r"""
\documentclass[final]{beamer}
\usepackage[orientation=portrait,size=a0,scale=1.35]{beamerposter}
\usetheme{default}
\usecolortheme{default}
\setbeamertemplate{navigation symbols}{}
\setbeamertemplate{headline}{}
\setbeamertemplate{footline}{}
\setbeamercolor{background canvas}{bg=bodyBg}

\usepackage{tcolorbox}
\tcbuselibrary{skins,fitting}
\usepackage{booktabs}
\usepackage{graphicx}
\usepackage{xcolor}
\usepackage{lmodern}
\usepackage{array}
\usepackage{ragged2e}
\usepackage{tabularx}
\usepackage{setspace}
\usepackage{microtype}
\usepackage{amsmath}

%% ── Colors ──────────────────────────────────────────────────────────────────
\definecolor{darkBlue}{RGB}{26,58,107}
\definecolor{midBlue}{RGB}{46,109,180}
\definecolor{lightBlue}{RGB}{232,242,255}
\definecolor{accent}{RGB}{240,165,0}
\definecolor{darkText}{RGB}{26,26,46}
\definecolor{bodyBg}{RGB}{245,247,250}
\definecolor{statLabel}{RGB}{200,220,248}
\definecolor{tblHdr}{RGB}{46,109,180}
\definecolor{tblAlt}{RGB}{240,245,255}
\definecolor{captionCol}{RGB}{68,85,119}

%% ── tcolorbox styles ────────────────────────────────────────────────────────
\tcbset{
  sectionblock/.style={
    enhanced,
    colback=lightBlue,
    colframe=midBlue,
    colbacktitle=darkBlue,
    coltitle=white,
    fonttitle=\bfseries\large,
    boxrule=1.5pt,
    arc=4pt,
    top=4pt, bottom=6pt, left=6pt, right=6pt,
    toptitle=3pt, bottomtitle=3pt,
    title={#1},
  },
  visualblock/.style={
    enhanced,
    colback=lightBlue,
    colframe=midBlue,
    colbacktitle=darkBlue,
    coltitle=white,
    fonttitle=\bfseries\normalsize,
    boxrule=1.5pt,
    arc=4pt,
    top=4pt, bottom=6pt, left=4pt, right=4pt,
    toptitle=3pt, bottomtitle=3pt,
    title={#1},
  },
}

%% ── Bullet helper ───────────────────────────────────────────────────────────
\newcommand{\pbullet}[1]{%
  \vspace{2pt}%
  \noindent$\blacktriangleright$\hspace{6pt}\begin{minipage}[t]{0.93\linewidth}%
  \RaggedRight #1\end{minipage}\par%
}
"""


def build_header(content: dict) -> str:
    title  = esc(content.get("title", ""))
    authors = esc(content.get("authors", ""))
    affils  = esc(content.get("affiliations", ""))
    conf    = esc(content.get("conference", ""))
    affconf = f"{affils} \\quad\\textbullet\\quad {conf}" if affils and conf else affils or conf

    return rf"""
%% ── HEADER ───────────────────────────────────────────────────────────────────
\begin{{beamercolorbox}}[wd=\paperwidth,sep=0pt]{{}}
  \begin{{tcolorbox}}[
    enhanced, colback=darkBlue, colframe=darkBlue,
    boxrule=0pt, arc=0pt, left=20pt, right=20pt, top=18pt, bottom=14pt,
    width=\paperwidth,
  ]
    \centering
    {{\bfseries\Huge\color{{white}} {title}}}\\[10pt]
    {{\large\color{{white}} {authors}}}\\[4pt]
    {{\normalsize\color{{statLabel}} {affconf}}}
  \end{{tcolorbox}}
\end{{beamercolorbox}}
"""


def build_banner(content: dict) -> str:
    takeaway = esc(content.get("key_takeaway", ""))
    return rf"""
%% ── TAKEAWAY BANNER ─────────────────────────────────────────────────────────
\begin{{tcolorbox}}[
  enhanced, colback=accent, colframe=accent,
  boxrule=0pt, arc=0pt, left=16pt, right=16pt, top=10pt, bottom=10pt,
  width=\paperwidth,
]
  \centering{{\bfseries\Large\color{{white}} $\star$\quad {takeaway}\quad$\star$}}
\end{{tcolorbox}}
\vspace{{4pt}}
"""


def build_stats_strip(key_numbers: list) -> str:
    if not key_numbers:
        return ""
    n = len(key_numbers)
    cell_pct = 0.98 / n

    cells = ""
    for i, item in enumerate(key_numbers):
        stat  = esc(item.get("stat", ""))
        label = esc(item.get("label", ""))
        sep = r"\hspace{2pt}\textcolor{white}{\rule[-8pt]{1pt}{36pt}}\hspace{2pt}" if i > 0 else ""
        cells += rf"""
    {sep}\begin{{minipage}}[c]{{{cell_pct:.3f}\linewidth}}
      \centering
      {{\bfseries\huge\color{{white}} {stat}}}\\[2pt]
      {{\small\color{{statLabel}} {label}}}
    \end{{minipage}}%"""

    return rf"""
%% ── KEY NUMBERS STRIP ───────────────────────────────────────────────────────
\begin{{tcolorbox}}[
  enhanced, colback=midBlue, colframe=midBlue,
  boxrule=0pt, arc=0pt, left=8pt, right=8pt, top=10pt, bottom=10pt,
  width=\paperwidth,
]
  \centering
  {cells}
\end{{tcolorbox}}
\vspace{{6pt}}
"""


def build_section(heading: str, bullets: list) -> str:
    bullet_lines = "\n".join(
        rf"    \pbullet{{{esc(b)}}}" for b in bullets
    )
    return rf"""
  \begin{{tcolorbox}}[sectionblock={{{esc(heading)}}}]
    \setstretch{{1.1}}
{bullet_lines}
  \end{{tcolorbox}}
  \vspace{{8pt}}
"""


def build_figure(caption: str, img_path: str) -> str:
    img_path_esc = img_path.replace("\\", "/")
    return rf"""
  \begin{{tcolorbox}}[visualblock={{{esc(caption)}}}]
    \centering
    \includegraphics[width=\linewidth,keepaspectratio]{{{img_path_esc}}}
  \end{{tcolorbox}}
  \vspace{{8pt}}
"""


def build_table(caption: str, headers: list, rows: list) -> str:
    n_cols = len(headers)
    if n_cols == 0:
        return ""

    # Column spec: first col left-aligned, rest centered
    col_spec = "l" + "c" * (n_cols - 1)

    hdr_row = " & ".join(
        rf"\textbf{{\color{{white}}{esc(h)}}}" for h in headers
    ) + r" \\"

    data_rows = []
    for i, row in enumerate(rows):
        cells = " & ".join(esc(str(v)) for v in row[:n_cols])
        bg = r"\rowcolor{tblAlt}" if i % 2 == 0 else ""
        data_rows.append(f"    {bg}{cells} \\\\")
    data_str = "\n".join(data_rows)

    return rf"""
  \begin{{tcolorbox}}[visualblock={{{esc(caption)}}}]
    \centering
    \setlength{{\tabcolsep}}{{6pt}}
    \renewcommand{{\arraystretch}}{{1.2}}
    \small
    \begin{{tabular}}{{{col_spec}}}
      \toprule
      \rowcolor{{tblHdr}} {hdr_row}
      \midrule
{data_str}
      \bottomrule
    \end{{tabular}}
  \end{{tcolorbox}}
  \vspace{{8pt}}
"""


def build_body(content: dict, figure_paths: list[str]) -> str:
    sections = content.get("sections", [])
    n = len(sections)
    left_secs = sections[: math.ceil(n / 2)]
    mid_secs  = sections[math.ceil(n / 2):]

    left_col  = "".join(build_section(s["heading"], s.get("bullets", [])) for s in left_secs)
    mid_col   = "".join(build_section(s["heading"], s.get("bullets", [])) for s in mid_secs)

    right_col = ""
    for fig, img_path in zip(content.get("key_figures", []), figure_paths):
        right_col += build_figure(fig.get("caption", "Figure"), img_path)
    for tbl in content.get("key_tables", []):
        right_col += build_table(
            tbl.get("caption", "Table"),
            tbl.get("headers", []),
            tbl.get("rows", []),
        )

    return rf"""
%% ── BODY ────────────────────────────────────────────────────────────────────
\begin{{columns}}[T,totalwidth=\linewidth]

  \begin{{column}}{{0.305\linewidth}}
{left_col}
  \end{{column}}
  \hspace{{4pt}}
  \begin{{column}}{{0.305\linewidth}}
{mid_col}
  \end{{column}}
  \hspace{{4pt}}
  \begin{{column}}{{0.36\linewidth}}
{right_col}
  \end{{column}}

\end{{columns}}
"""


def build_tex(content: dict, figure_paths: list[str]) -> str:
    return (
        LATEX_PREAMBLE
        + r"\begin{document}" + "\n"
        + r"\begin{frame}[t,fragile]" + "\n"
        + build_header(content)
        + build_banner(content)
        + build_stats_strip(content.get("key_numbers", []))
        + build_body(content, figure_paths)
        + r"\end{frame}" + "\n"
        + r"\end{document}" + "\n"
    )


# ---------------------------------------------------------------------------
# Compilation
# ---------------------------------------------------------------------------

def compile_tex(tex_path: Path) -> Path | None:
    """Run pdflatex twice (for correct layout) in the tex file's directory."""
    pdflatex = shutil.which("pdflatex") or "/Library/TeX/texbin/pdflatex"
    if not Path(pdflatex).exists():
        print("  [warn] pdflatex not found — skipping compilation.", file=sys.stderr)
        return None

    pdf_path = tex_path.with_suffix(".pdf")
    for run in range(2):
        result = subprocess.run(
            [pdflatex, "-interaction=nonstopmode", "-output-directory",
             str(tex_path.parent), str(tex_path)],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            # Print last 20 lines of log for diagnosis
            log = tex_path.with_suffix(".log")
            if log.exists():
                lines = log.read_text(errors="replace").splitlines()
                print("\n".join(lines[-20:]), file=sys.stderr)
            print(f"  [error] pdflatex failed (run {run+1}).", file=sys.stderr)
            return None

    return pdf_path if pdf_path.exists() else None


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Render poster_content.json to a LaTeX/PDF poster.")
    parser.add_argument("content", help="Path to poster_content.json")
    parser.add_argument("--slides",     help="Slides .pptx for figure extraction")
    parser.add_argument("--out",        help="Output .tex path")
    parser.add_argument("--no-compile", action="store_true",
                        help="Generate .tex only, skip pdflatex")
    args = parser.parse_args()

    content_path = Path(args.content)
    if not content_path.exists():
        print(f"Error: file not found: {content_path}", file=sys.stderr)
        sys.exit(1)

    out_tex = (
        Path(args.out) if args.out
        else content_path.with_stem(
            content_path.stem.replace("_poster_content", "") + "_poster"
        ).with_suffix(".tex")
    )

    content = json.loads(content_path.read_text())
    print(f"Building poster: {content.get('title', '')[:60]}...")

    # ── Extract figures ────────────────────────────────────────────────────
    figure_paths: list[str] = []
    slides_path = args.slides

    for fig in content.get("key_figures", []):
        # Manual override takes priority
        image_path = fig.get("image_path")
        if image_path and Path(image_path).exists():
            figure_paths.append(str(Path(image_path).resolve()))
            print(f"  Figure: using {image_path}")
            continue

        if not slides_path:
            print(f"  [info] Skipping figure (slide {fig['slide']}) — no --slides provided.")
            continue

        # Save figure next to the .tex file
        fig_out = out_tex.parent / f"fig_slide{fig['slide']}.png"
        print(f"  Extracting figure from slide {fig['slide']}...")
        if extract_slide_image(slides_path, fig["slide"], fig_out):
            figure_paths.append(str(fig_out.resolve()))
        else:
            print(f"  [warn] Could not extract slide {fig['slide']} — figure skipped.",
                  file=sys.stderr)

    # ── Generate .tex ──────────────────────────────────────────────────────
    tex_src = build_tex(content, figure_paths)
    out_tex.write_text(tex_src, encoding="utf-8")
    print(f"Saved:  {out_tex}")

    # ── Compile ────────────────────────────────────────────────────────────
    if args.no_compile:
        return

    print("Compiling with pdflatex (2 passes)...")
    pdf = compile_tex(out_tex)
    if pdf:
        print(f"Saved:  {pdf}")
    else:
        print("Compilation failed — check the .log file.", file=sys.stderr)


if __name__ == "__main__":
    main()
