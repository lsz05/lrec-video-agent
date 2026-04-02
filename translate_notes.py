"""
translate_notes.py

Translates Japanese speaker notes in a PowerPoint (.pptx) into English using the Claude API,
then writes the translated notes into the target (already slide-translated) .pptx.

Usage:
    python translate_notes.py <source.pptx> <target.pptx> [--paper <paper.pdf>]

    <source.pptx>   Original Japanese .pptx (notes are read from here)
    <target.pptx>   English slide .pptx produced by translate_slides.py (notes are written here)
    --paper         Optional PDF of the paper for terminology context

The target file is updated in-place.

Requirements:
    pip install anthropic python-pptx PyMuPDF
    export ANTHROPIC_API_KEY=...
"""

import argparse
import json
import os
import sys
import time
from pathlib import Path

import anthropic
from pptx import Presentation
from pptx.util import Pt
from lxml import etree


# ---------------------------------------------------------------------------
# PDF context (reused from translate_slides.py)
# ---------------------------------------------------------------------------

def extract_pdf_abstract(pdf_path: str, max_chars: int = 3000) -> str:
    try:
        import fitz
    except ImportError:
        print("  [warn] PyMuPDF not installed — skipping paper context.", file=sys.stderr)
        return ""
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) >= max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        print(f"  [warn] Could not read PDF ({pdf_path}): {e}", file=sys.stderr)
        return ""


# ---------------------------------------------------------------------------
# Notes extraction
# ---------------------------------------------------------------------------

def extract_notes(prs: Presentation) -> list[str]:
    """Return a list of speaker note strings, one per slide (empty string if none)."""
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            text = ""
        notes.append(text)
    return notes


# ---------------------------------------------------------------------------
# Translation via Claude
# ---------------------------------------------------------------------------

SYSTEM_PROMPT = """\
You are a professional academic translator specializing in NLP / AI research.
You will receive the speaker notes for one slide of a Japanese academic presentation.
Translate the notes into natural, fluent spoken English suitable for delivery at an international conference.

Rules:
- Translate faithfully — preserve all content and nuance.
- Keep the spoken, presenter-facing register (first person, natural flow).
- Keep technical terms, model names, dataset names, and acronyms unchanged.
- Keep numbers, formulas, and citations unchanged.
- Do NOT add, remove, or significantly expand the content.
- Return ONLY the translated text, no explanations or markdown.
"""

MAX_RETRIES = 3
RETRY_BACKOFF = [2, 5, 15]


def translate_note(
    client: anthropic.Anthropic,
    note: str,
    slide_number: int,
    paper_context: str = "",
) -> str:
    """Translate one slide's speaker notes. Returns original on total failure."""
    if not note.strip():
        return ""

    context_block = (
        f"\nPaper context (for terminology reference):\n{paper_context}\n\n"
        if paper_context else ""
    )
    user_content = f"Slide {slide_number}:\n{context_block}{note}"

    last_error = None
    for attempt in range(MAX_RETRIES):
        try:
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2048,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": user_content}],
            )
            return message.content[0].text.strip()

        except anthropic.RateLimitError as e:
            last_error = f"Rate limit: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except anthropic.APIStatusError as e:
            last_error = f"API error {e.status_code}: {e.message}"
            if e.status_code < 500:
                break
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except anthropic.APIConnectionError as e:
            last_error = f"Connection error: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

    print(f"\n    [error] Slide {slide_number}: all retries failed ({last_error}). "
          "Keeping original text.", file=sys.stderr)
    return note  # fall back to Japanese original


# ---------------------------------------------------------------------------
# Writing notes back into the target PPTX
# ---------------------------------------------------------------------------

def set_notes(slide, text: str) -> None:
    """Set the speaker notes text for a slide, creating the notes slide if needed."""
    notes_slide = slide.notes_slide  # creates one if it doesn't exist
    tf = notes_slide.notes_text_frame

    # Clear all existing paragraphs except the first
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)

    # Write into the first paragraph (split on newlines to preserve paragraph breaks)
    lines = text.split("\n")
    tf.paragraphs[0].runs[0].text = lines[0] if lines else ""

    # Add additional paragraphs for each subsequent line
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


# ---------------------------------------------------------------------------
# Checkpoint helpers
# ---------------------------------------------------------------------------

def checkpoint_path(target_path: Path) -> Path:
    return target_path.with_suffix(".notes_checkpoint.json")


def load_checkpoint(target_path: Path) -> dict[int, str] | None:
    cp = checkpoint_path(target_path)
    if cp.exists():
        try:
            data = json.loads(cp.read_text())
            print(f"  Resuming from checkpoint ({len(data)} notes already translated)")
            return {int(k): v for k, v in data.items()}
        except Exception as e:
            print(f"  [warn] Could not read checkpoint: {e} — starting fresh", file=sys.stderr)
    return None


def save_checkpoint(target_path: Path, done: dict[int, str]) -> None:
    checkpoint_path(target_path).write_text(json.dumps(done, ensure_ascii=False, indent=2))


def remove_checkpoint(target_path: Path) -> None:
    cp = checkpoint_path(target_path)
    if cp.exists():
        cp.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Translate Japanese speaker notes into English and write into target PPTX.")
    parser.add_argument("source", help="Source .pptx with Japanese speaker notes")
    parser.add_argument("target", help="Target .pptx to write English notes into (modified in-place)")
    parser.add_argument("--paper", help="Optional paper PDF for terminology context")
    args = parser.parse_args()

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY is not set.\n"
              "  export ANTHROPIC_API_KEY=sk-ant-...", file=sys.stderr)
        sys.exit(1)

    source_path = Path(args.source)
    target_path = Path(args.target)

    for p in [source_path, target_path]:
        if not p.exists():
            print(f"Error: file not found: {p}", file=sys.stderr)
            sys.exit(1)

    paper_context = ""
    if args.paper:
        paper_pdf = Path(args.paper)
        if not paper_pdf.exists():
            print(f"  [warn] Paper PDF not found: {paper_pdf} — skipping context",
                  file=sys.stderr)
        else:
            print(f"Extracting paper context from {paper_pdf.name}...")
            paper_context = extract_pdf_abstract(str(paper_pdf))

    client = anthropic.Anthropic()

    print(f"Loading source: {source_path.name}")
    src_prs = Presentation(str(source_path))
    japanese_notes = extract_notes(src_prs)
    total = len(japanese_notes)

    notes_with_text = sum(1 for n in japanese_notes if n.strip())
    print(f"Found {notes_with_text}/{total} slides with notes.\n")

    checkpoint = load_checkpoint(target_path) or {}

    print(f"Translating notes...\n")
    translated_notes: list[str] = []
    try:
        for i, note in enumerate(japanese_notes, start=1):
            if not note.strip():
                print(f"  Slide {i:02d}/{total}  (no notes — skipped)")
                translated_notes.append("")
                continue

            if i in checkpoint:
                print(f"  Slide {i:02d}/{total}  (from checkpoint)")
                translated_notes.append(checkpoint[i])
                continue

            print(f"  Slide {i:02d}/{total}  ({len(note)} chars)...", end=" ", flush=True)
            translated = translate_note(client, note, i, paper_context)
            translated_notes.append(translated)
            checkpoint[i] = translated
            save_checkpoint(target_path, checkpoint)
            print("done")

    except KeyboardInterrupt:
        print("\n\nInterrupted. Progress saved — re-run to resume.", file=sys.stderr)
        sys.exit(1)

    print(f"\nWriting translated notes into {target_path.name}...")
    tgt_prs = Presentation(str(target_path))

    if len(tgt_prs.slides) != len(translated_notes):
        print(f"  [warn] Slide count mismatch: source has {total}, "
              f"target has {len(tgt_prs.slides)}. Writing what we can.", file=sys.stderr)

    for i, (slide, note) in enumerate(zip(tgt_prs.slides, translated_notes), start=1):
        if note.strip():
            try:
                set_notes(slide, note)
            except Exception as e:
                print(f"  [warn] Slide {i}: could not write notes: {e}", file=sys.stderr)

    tgt_prs.save(str(target_path))
    remove_checkpoint(target_path)
    print(f"\nSaved: {target_path}")


if __name__ == "__main__":
    main()
