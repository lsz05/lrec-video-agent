"""
render_poster.py

Renders a poster_content.json file into an A0-sized academic poster PPTX.

Layout (portrait A0, 33.11" × 46.81"):
  ┌──────────────────────────────────────────┐
  │  HEADER  title / authors / affiliations  │  ~5.5"
  ├──────────────────────────────────────────┤
  │  KEY TAKEAWAY (banner quote)             │  ~1.5"
  ├──────────────────────────────────────────┤
  │  KEY NUMBERS STRIP  28 | 5 | 40 | >0.99 │  ~2.2"
  ├───────────────┬───────────────┬──────────┤
  │  LEFT         │  MIDDLE       │  RIGHT   │  ~35"
  │  sections 1–3 │  sections 4–6 │  figures │
  │               │               │  tables  │
  └───────────────┴───────────────┴──────────┘

Usage:
    python render_poster.py <poster_content.json> [--slides <slides_en.pptx>] [--out <poster.pptx>]

    --slides  English slides .pptx used to extract figure images (optional; figures skipped if omitted)

Requirements:
    pip install python-pptx
"""

import argparse
import io
import json
import math
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------

SLIDE_W = 33.11   # A0 portrait width  (inches)
SLIDE_H = 46.81   # A0 portrait height (inches)

MARGIN = 0.4
GUTTER = 0.3

# Three-column widths: left/middle share equally; right is slightly wider for visuals
TEXT_COLS  = 2
VISUAL_COL_FRAC = 0.33   # right column takes 33% of usable width
USABLE_W   = SLIDE_W - MARGIN * 2 - GUTTER * (TEXT_COLS)   # 2 gutters for 3 cols
VISUAL_W   = USABLE_W * VISUAL_COL_FRAC
TEXT_COL_W = (USABLE_W - VISUAL_W) / TEXT_COLS

LEFT_X   = MARGIN
MID_X    = MARGIN + TEXT_COL_W + GUTTER
RIGHT_X  = MARGIN + TEXT_COL_W * 2 + GUTTER * 2

# Vertical zones
HEADER_TOP = MARGIN
HEADER_H   = 5.5
BANNER_TOP = HEADER_TOP + HEADER_H + 0.2
BANNER_H   = 1.5
STATS_TOP  = BANNER_TOP + BANNER_H + 0.2
STATS_H    = 2.2
BODY_TOP   = STATS_TOP + STATS_H + 0.3
BODY_H     = SLIDE_H - BODY_TOP - MARGIN

SECTION_GAP = 0.25
VISUAL_GAP  = 0.3

# Fonts (large — A0 is read from ~1m away)
FONT_TITLE    = 64
FONT_AUTHORS  = 34
FONT_CONF     = 28
FONT_BANNER   = 34
FONT_STAT     = 52
FONT_STAT_LB  = 22
FONT_HEADING  = 38
FONT_BULLET   = 25
FONT_TBL_HDR  = 22
FONT_TBL_CELL = 20
FONT_CAPTION  = 21

# Colors
C_DARK_BLUE  = RGBColor(0x1a, 0x3a, 0x6b)
C_MID_BLUE   = RGBColor(0x2e, 0x6d, 0xb4)
C_LIGHT_BLUE = RGBColor(0xe8, 0xf2, 0xff)
C_TBL_HDR    = RGBColor(0x2e, 0x6d, 0xb4)
C_TBL_ALT    = RGBColor(0xf0, 0xf5, 0xff)
C_ACCENT     = RGBColor(0xf0, 0xa5, 0x00)
C_WHITE      = RGBColor(0xff, 0xff, 0xff)
C_DARK_TEXT  = RGBColor(0x1a, 0x1a, 0x2e)
C_BODY_BG    = RGBColor(0xf5, 0xf7, 0xfa)
C_CAPTION    = RGBColor(0x44, 0x55, 0x77)


# ---------------------------------------------------------------------------
# Generic shape helpers
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_label(slide, left, top, width, height, text, font_size,
              bold=False, color=None, bg=None, align=PP_ALIGN.LEFT, word_wrap=True):
    if bg:
        shape = add_rect(slide, left, top, width, height, fill_color=bg)
        tf = shape.text_frame
    else:
        tf = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


# ---------------------------------------------------------------------------
# Key numbers strip
# ---------------------------------------------------------------------------

def render_stats_strip(slide, key_numbers: list[dict]):
    if not key_numbers:
        return
    n = len(key_numbers)
    cell_w = (SLIDE_W - MARGIN * 2) / n
    add_rect(slide, 0, STATS_TOP, SLIDE_W, STATS_H, fill_color=C_MID_BLUE)
    for i, item in enumerate(key_numbers):
        x = MARGIN + i * cell_w
        if i > 0:
            add_rect(slide, x - 0.02, STATS_TOP + 0.25, 0.03, STATS_H - 0.5,
                     fill_color=C_WHITE)
        add_label(slide, x + 0.1, STATS_TOP + 0.15, cell_w - 0.2, 1.3,
                  str(item.get("stat", "")), FONT_STAT, bold=True,
                  color=C_WHITE, align=PP_ALIGN.CENTER)
        add_label(slide, x + 0.1, STATS_TOP + 1.4, cell_w - 0.2, 0.7,
                  str(item.get("label", "")), FONT_STAT_LB,
                  color=RGBColor(0xc8, 0xdc, 0xf8), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Text section box
# ---------------------------------------------------------------------------

def render_section(slide, heading, bullets, left, top, width, max_height) -> float:
    HEADING_H = 0.62
    BULLET_H  = 0.43
    PADDING   = 0.14

    box_h = min(HEADING_H + len(bullets) * BULLET_H + PADDING * 2, max_height)

    add_rect(slide, left, top, width, box_h, fill_color=C_LIGHT_BLUE, line_color=C_MID_BLUE)
    add_rect(slide, left, top, width, HEADING_H, fill_color=C_DARK_BLUE)
    add_label(slide, left + 0.12, top + 0.07, width - 0.2, HEADING_H - 0.1,
              heading, FONT_HEADING, bold=True, color=C_WHITE)

    tb = slide.shapes.add_textbox(
        Inches(left + 0.18), Inches(top + HEADING_H + PADDING),
        Inches(width - 0.32), Inches(box_h - HEADING_H - PADDING * 2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "▸  " + bullet
        run.font.size = Pt(FONT_BULLET)
        run.font.color.rgb = C_DARK_TEXT

    return box_h


# ---------------------------------------------------------------------------
# Table box
# ---------------------------------------------------------------------------

def render_table_box(slide, caption: str, headers: list, rows: list,
                     left: float, top: float, width: float, max_height: float) -> float:
    """Render a native PPTX table with caption. Returns height used."""
    CAPTION_H = 0.5
    HDR_H     = 0.45
    ROW_H     = 0.38
    PADDING   = 0.15

    n_rows = len(rows)
    n_cols = len(headers)
    table_h = HDR_H + n_rows * ROW_H
    box_h = min(CAPTION_H + table_h + PADDING * 2, max_height)
    table_h = box_h - CAPTION_H - PADDING * 2

    # Outer box
    add_rect(slide, left, top, width, box_h, fill_color=C_LIGHT_BLUE, line_color=C_MID_BLUE)

    # Caption bar
    add_rect(slide, left, top, width, CAPTION_H, fill_color=C_DARK_BLUE)
    add_label(slide, left + 0.12, top + 0.06, width - 0.2, CAPTION_H - 0.1,
              caption, FONT_CAPTION, bold=True, color=C_WHITE)

    # Table
    tbl_top  = top + CAPTION_H + PADDING
    tbl_left = left + PADDING

    # Distribute column widths evenly
    tbl_w = width - PADDING * 2
    shape = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(tbl_left), Inches(tbl_top),
        Inches(tbl_w), Inches(table_h))
    table = shape.table

    col_w = tbl_w / n_cols
    for ci in range(n_cols):
        table.columns[ci].width = Inches(col_w)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TBL_HDR
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(hdr)
        run.font.size = Pt(FONT_TBL_HDR)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_TBL_ALT if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(FONT_TBL_CELL)
            run.font.color.rgb = C_DARK_TEXT

    return box_h


# ---------------------------------------------------------------------------
# Figure extraction from slides PPTX
# ---------------------------------------------------------------------------

def extract_slide_image(pptx_path: str, slide_number: int) -> bytes | None:
    """
    Extract the largest embedded picture from a slide and return its raw bytes.
    Falls back to the second-largest if the largest is a background/logo.
    Returns None if no picture shapes are found.
    """
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(pptx_path)
    except Exception as e:
        print(f"  [warn] Could not open slides file: {e}", file=sys.stderr)
        return None

    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"  [warn] Slide {slide_number} out of range ({len(prs.slides)} slides).",
              file=sys.stderr)
        return None

    slide = prs.slides[slide_number - 1]

    # Collect all picture shapes sorted by area (largest first)
    pictures = [
        shape for shape in slide.shapes
        if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    ]
    if not pictures:
        print(f"  [warn] Slide {slide_number} has no embedded pictures.", file=sys.stderr)
        return None

    pictures.sort(key=lambda s: s.width * s.height, reverse=True)
    for pic in pictures:
        try:
            blob = pic.image.blob
            if blob:
                return blob
        except (ValueError, AttributeError):
            continue  # linked or shape-only picture — try next

    print(f"  [warn] Slide {slide_number} pictures have no embedded image data "
          f"(may be drawn shapes or linked). Skipping.", file=sys.stderr)
    return None


def render_figure_box(slide, img_bytes: bytes, caption: str,
                      left: float, top: float, width: float, max_height: float) -> float:
    """Insert a figure image with caption. Returns height used."""
    CAPTION_H = 0.5
    PADDING   = 0.1
    img_h     = min(max_height - CAPTION_H - PADDING * 2, 8.5)
    box_h     = CAPTION_H + img_h + PADDING * 2

    # Caption bar
    add_rect(slide, left, top, width, CAPTION_H, fill_color=C_DARK_BLUE)
    add_label(slide, left + 0.12, top + 0.06, width - 0.2, CAPTION_H - 0.1,
              caption, FONT_CAPTION, bold=True, color=C_WHITE)

    # Image
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(img_bytes)
        tmp_path = f.name

    try:
        slide.shapes.add_picture(
            tmp_path,
            Inches(left + PADDING),
            Inches(top + CAPTION_H + PADDING),
            width=Inches(width - PADDING * 2))
    except Exception as e:
        print(f"  [warn] Could not insert figure image: {e}", file=sys.stderr)
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    # Outline box
    add_rect(slide, left, top, width, box_h, line_color=C_MID_BLUE)

    return box_h


# ---------------------------------------------------------------------------
# Main renderer
# ---------------------------------------------------------------------------

def render(content: dict, out_path: Path, slides_path: str | None = None):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_BODY_BG)

    # ── HEADER ───────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, SLIDE_W, HEADER_TOP + HEADER_H, fill_color=C_DARK_BLUE)
    add_label(slide, MARGIN, HEADER_TOP + 0.3, SLIDE_W - MARGIN * 2, 3.2,
              content.get("title", ""), FONT_TITLE, bold=True,
              color=C_WHITE, align=PP_ALIGN.CENTER)
    add_label(slide, MARGIN, HEADER_TOP + 3.6, SLIDE_W - MARGIN * 2, 1.0,
              content.get("authors", ""), FONT_AUTHORS,
              color=C_WHITE, align=PP_ALIGN.CENTER)
    affil = content.get("affiliations", "")
    conf  = content.get("conference", "")
    add_label(slide, MARGIN, HEADER_TOP + 4.55, SLIDE_W - MARGIN * 2, 0.7,
              f"{affil}    ·    {conf}" if affil and conf else affil or conf,
              FONT_CONF, color=RGBColor(0xb0, 0xc8, 0xf0), align=PP_ALIGN.CENTER)

    # ── TAKEAWAY BANNER ──────────────────────────────────────────────────────
    add_rect(slide, 0, BANNER_TOP, SLIDE_W, BANNER_H, fill_color=C_ACCENT)
    add_label(slide, MARGIN, BANNER_TOP + 0.15, SLIDE_W - MARGIN * 2, BANNER_H - 0.2,
              "✦  " + content.get("key_takeaway", ""),
              FONT_BANNER, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── KEY NUMBERS STRIP ────────────────────────────────────────────────────
    render_stats_strip(slide, content.get("key_numbers", []))

    # ── BODY: 3-COLUMN ───────────────────────────────────────────────────────
    sections = content.get("sections", [])
    n = len(sections)
    left_secs  = sections[: math.ceil(n / 2)]
    mid_secs   = sections[math.ceil(n / 2) :]

    def render_text_column(sec_list, col_x):
        y = BODY_TOP
        rem = BODY_H
        for sec in sec_list:
            if rem < 1.0:
                break
            used = render_section(slide, sec["heading"], sec.get("bullets", []),
                                  col_x, y, TEXT_COL_W, rem)
            y += used + SECTION_GAP
            rem -= used + SECTION_GAP

    render_text_column(left_secs, LEFT_X)
    render_text_column(mid_secs,  MID_X)

    # ── RIGHT COLUMN: FIGURES + TABLES ───────────────────────────────────────
    y = BODY_TOP
    rem = BODY_H

    # Figures
    for fig in content.get("key_figures", []):
        if rem < 2.0:
            break

        img_bytes = None

        # 1. Manual override — image_path in the JSON takes priority
        image_path = fig.get("image_path")
        if image_path:
            p = Path(image_path)
            if p.exists():
                img_bytes = p.read_bytes()
                print(f"  Using provided image: {p.name}")
            else:
                print(f"  [warn] image_path not found: {image_path}", file=sys.stderr)

        # 2. Extract from slides PPTX
        if img_bytes is None:
            if not slides_path:
                print(f"  [info] Skipping figure (slide {fig['slide']}) — no --slides provided.")
                continue
            print(f"  Extracting figure from slide {fig['slide']}...")
            img_bytes = extract_slide_image(slides_path, fig["slide"])

        if img_bytes:
            used = render_figure_box(slide, img_bytes,
                                     fig.get("caption", "Figure")[:70],
                                     RIGHT_X, y, VISUAL_W, rem)
            y += used + VISUAL_GAP
            rem -= used + VISUAL_GAP

    # Tables
    for tbl in content.get("key_tables", []):
        if rem < 2.0:
            break
        used = render_table_box(slide,
                                tbl.get("caption", "Table"),
                                tbl.get("headers", []),
                                tbl.get("rows", []),
                                RIGHT_X, y, VISUAL_W, rem)
        y += used + VISUAL_GAP
        rem -= used + VISUAL_GAP

    prs.save(str(out_path))
    print(f"Saved: {out_path}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Render poster_content.json into an A0 PPTX poster.")
    parser.add_argument("content", help="Path to poster_content.json")
    parser.add_argument("--slides", help="English slides .pptx for figure extraction (optional)")
    parser.add_argument("--out",    help="Output .pptx path")
    args = parser.parse_args()

    content_path = Path(args.content)
    if not content_path.exists():
        print(f"Error: file not found: {content_path}", file=sys.stderr)
        sys.exit(1)

    out_path = (
        Path(args.out) if args.out
        else content_path.with_stem(
            content_path.stem.replace("_poster_content", "") + "_poster"
        ).with_suffix(".pptx")
    )

    slides_path = args.slides
    if slides_path and not Path(slides_path).exists():
        print(f"  [warn] Slides file not found: {slides_path} — figures will be skipped.",
              file=sys.stderr)
        slides_path = None

    content = json.loads(content_path.read_text())
    print(f"Rendering poster: {content.get('title', '')[:60]}...")
    render(content, out_path, slides_path)


if __name__ == "__main__":
    main()
