"""
tts_notes.py

Reads English speaker notes from a .pptx and synthesizes each slide's notes
into an audio file using the OpenAI TTS API.

Usage:
    python tts_notes.py <slides_en.pptx> [--out-dir <dir>] [--voice <voice>] [--model <model>]

    <slides_en.pptx>   English .pptx produced by translate_notes.py
    --out-dir          Directory to write audio files into (default: <pptx_stem>_audio/)
    --voice            OpenAI TTS voice (default: alloy)
                       Choices: alloy, echo, fable, onyx, nova, shimmer
    --model            OpenAI TTS model (default: tts-1-hd)
                       Choices: tts-1, tts-1-hd

Output:
    One .mp3 per slide: slide_01.mp3, slide_02.mp3, ...
    Slides with no notes produce no file.
    A manifest.json is written listing each slide's audio file and duration.

Requirements:
    pip install openai python-pptx mutagen
    export OPENAI_API_KEY=sk-...
"""

import argparse
import json
import os
import sys
import time
from pathlib import Path

import openai
from pptx import Presentation


# ---------------------------------------------------------------------------
# Notes extraction
# ---------------------------------------------------------------------------

def extract_notes(prs: Presentation) -> list[str]:
    """Return a list of speaker note strings, one per slide (empty string if none)."""
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            text = ""
        notes.append(text)
    return notes


# ---------------------------------------------------------------------------
# Audio duration helper
# ---------------------------------------------------------------------------

def get_mp3_duration(path: Path) -> float:
    """Return duration in seconds of an MP3 file using mutagen."""
    try:
        from mutagen.mp3 import MP3
        return MP3(str(path)).info.length
    except Exception:
        return 0.0


# ---------------------------------------------------------------------------
# TTS via OpenAI (with retry)
# ---------------------------------------------------------------------------

MAX_RETRIES = 3
RETRY_BACKOFF = [2, 5, 15]


def synthesize(
    client: openai.OpenAI,
    text: str,
    out_path: Path,
    voice: str,
    model: str,
    slide_number: int,
) -> bool:
    """
    Synthesize text to speech and save as MP3 at out_path.
    Returns True on success, False on total failure.
    """
    last_error = None
    for attempt in range(MAX_RETRIES):
        try:
            response = client.audio.speech.create(
                model=model,
                voice=voice,
                input=text,
                response_format="mp3",
            )
            out_path.write_bytes(response.content)
            return True

        except openai.RateLimitError as e:
            last_error = f"Rate limit: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except openai.APIStatusError as e:
            last_error = f"API error {e.status_code}: {e.message}"
            if e.status_code < 500:
                break  # 4xx won't fix itself
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except openai.APIConnectionError as e:
            last_error = f"Connection error: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

    print(f"\n    [error] Slide {slide_number}: all retries failed ({last_error}). Skipping.",
          file=sys.stderr)
    return False


# ---------------------------------------------------------------------------
# Checkpoint helpers
# ---------------------------------------------------------------------------

def checkpoint_path(out_dir: Path) -> Path:
    return out_dir / "tts_checkpoint.json"


def load_checkpoint(out_dir: Path) -> set[int]:
    cp = checkpoint_path(out_dir)
    if cp.exists():
        try:
            done = set(json.loads(cp.read_text()))
            print(f"  Resuming from checkpoint ({len(done)} slides already synthesized)")
            return done
        except Exception as e:
            print(f"  [warn] Could not read checkpoint: {e} — starting fresh", file=sys.stderr)
    return set()


def save_checkpoint(out_dir: Path, done: set[int]) -> None:
    checkpoint_path(out_dir).write_text(json.dumps(sorted(done)))


def remove_checkpoint(out_dir: Path) -> None:
    cp = checkpoint_path(out_dir)
    if cp.exists():
        cp.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Synthesize English speaker notes to MP3 using OpenAI TTS.")
    parser.add_argument("input", help="Path to the English .pptx (with speaker notes)")
    parser.add_argument("--out-dir", help="Output directory for audio files")
    parser.add_argument("--voice", default="alloy",
                        choices=["alloy", "echo", "fable", "onyx", "nova", "shimmer"],
                        help="TTS voice (default: alloy)")
    parser.add_argument("--model", default="tts-1-hd",
                        choices=["tts-1", "tts-1-hd"],
                        help="TTS model (default: tts-1-hd)")
    args = parser.parse_args()

    if not os.environ.get("OPENAI_API_KEY"):
        print("Error: OPENAI_API_KEY is not set.\n"
              "  export OPENAI_API_KEY=sk-...", file=sys.stderr)
        sys.exit(1)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    out_dir = Path(args.out_dir) if args.out_dir else input_path.parent / (input_path.stem + "_audio")
    out_dir.mkdir(parents=True, exist_ok=True)

    client = openai.OpenAI()  # reads OPENAI_API_KEY from env

    print(f"Loading {input_path.name}...")
    prs = Presentation(str(input_path))
    notes = extract_notes(prs)
    total = len(notes)
    notes_with_text = sum(1 for n in notes if n.strip())
    print(f"Found {notes_with_text}/{total} slides with notes.")
    print(f"Voice: {args.voice}  Model: {args.model}  Output: {out_dir}/\n")

    done_slides = load_checkpoint(out_dir)
    manifest = []

    try:
        for i, note in enumerate(notes, start=1):
            entry = {"slide": i, "audio_file": None, "duration_sec": 0.0, "text": note}

            if not note.strip():
                print(f"  Slide {i:02d}/{total}  (no notes — skipped)")
                manifest.append(entry)
                continue

            out_file = out_dir / f"slide_{i:02d}.mp3"

            if i in done_slides and out_file.exists():
                print(f"  Slide {i:02d}/{total}  (from checkpoint)")
                entry["audio_file"] = str(out_file)
                entry["duration_sec"] = get_mp3_duration(out_file)
                manifest.append(entry)
                continue

            print(f"  Slide {i:02d}/{total}  ({len(note)} chars)...", end=" ", flush=True)
            success = synthesize(client, note, out_file, args.voice, args.model, i)

            if success:
                duration = get_mp3_duration(out_file)
                entry["audio_file"] = str(out_file)
                entry["duration_sec"] = duration
                done_slides.add(i)
                save_checkpoint(out_dir, done_slides)
                print(f"done  ({duration:.1f}s)")
            else:
                print("FAILED")

            manifest.append(entry)

    except KeyboardInterrupt:
        print("\n\nInterrupted. Progress saved — re-run to resume.", file=sys.stderr)
        _write_manifest(out_dir, manifest)
        sys.exit(1)

    _write_manifest(out_dir, manifest)
    remove_checkpoint(out_dir)

    total_duration = sum(e["duration_sec"] for e in manifest)
    print(f"\nTotal audio duration: {total_duration / 60:.1f} min  ({total_duration:.0f}s)")
    if total_duration > 0:
        target = 15 * 60
        diff = total_duration - target
        if abs(diff) > 30:
            direction = "over" if diff > 0 else "under"
            print(f"  [info] {abs(diff):.0f}s {direction} the 15-minute target — "
                  "consider trimming or expanding the notes.")
    print(f"Manifest: {out_dir / 'manifest.json'}")


def _write_manifest(out_dir: Path, manifest: list[dict]) -> None:
    (out_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
