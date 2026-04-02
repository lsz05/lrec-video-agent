"""
export_notes.py

Exports all speaker notes from a .pptx file to a plain text or JSON file.

Usage:
    python3 export_notes.py <slides.pptx> [--out <output>] [--format txt|json]

    <slides.pptx>   Any .pptx file with speaker notes
    --out           Output file path (default: <pptx_stem>_notes.txt or .json)
    --format        Output format: txt (default) or json

Text format:
    --- Slide 1 ---
    <note text>

    --- Slide 2 ---
    <note text>
    ...

JSON format:
    [{"slide": 1, "text": "..."}, ...]

Requirements:
    pip install python-pptx
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


def extract_notes(prs: Presentation) -> list[dict]:
    result = []
    for i, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
        else:
            text = ""
        if text:
            result.append({"slide": i, "text": text})
    return result


def main():
    parser = argparse.ArgumentParser(description="Export speaker notes from a .pptx file.")
    parser.add_argument("input", help="Path to the .pptx file")
    parser.add_argument("--out", help="Output file path")
    parser.add_argument("--format", choices=["txt", "json"], default="txt",
                        help="Output format: txt (default) or json")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    ext = "json" if args.format == "json" else "txt"
    out_path = Path(args.out) if args.out else input_path.with_stem(input_path.stem + "_notes").with_suffix(f".{ext}")

    prs = Presentation(str(input_path))
    notes = extract_notes(prs)

    if args.format == "json":
        out_path.write_text(json.dumps(notes, ensure_ascii=False, indent=2))
    else:
        lines = []
        for entry in notes:
            lines.append(f"--- Slide {entry['slide']} ---")
            lines.append(entry["text"])
            lines.append("")
        out_path.write_text("\n".join(lines), encoding="utf-8")

    total_words = sum(len(e["text"].split()) for e in notes)
    print(f"Exported {len(notes)} slides ({total_words} words) → {out_path}")


if __name__ == "__main__":
    main()
