"""
translate_slides.py

Translates a Japanese PowerPoint (.pptx) into English using the Claude API.
Preserves all formatting (fonts, sizes, colors, layout) — only text content is changed.

Usage:
    python translate_slides.py <input.pptx> [--output <output.pptx>] [--paper <paper.pdf>]

    <input.pptx>          Path to the input .pptx file
    --output              Output path (default: <input>_en.pptx next to input)
    --paper               Optional path to the paper PDF for terminology context

Requirements:
    pip install anthropic python-pptx PyMuPDF
    export ANTHROPIC_API_KEY=...
"""

import argparse
import json
import os
import sys
import time
from pathlib import Path

import anthropic
from pptx import Presentation


# ---------------------------------------------------------------------------
# PDF text extraction (optional, for paper context)
# ---------------------------------------------------------------------------

def extract_pdf_abstract(pdf_path: str, max_chars: int = 3000) -> str:
    """Extract the first ~max_chars chars from a PDF (typically abstract + intro)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("  [warn] PyMuPDF not installed — skipping paper context. Run: pip install PyMuPDF",
              file=sys.stderr)
        return ""
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) >= max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        print(f"  [warn] Could not read paper PDF ({pdf_path}): {e}", file=sys.stderr)
        return ""


# ---------------------------------------------------------------------------
# Slide text extraction
# ---------------------------------------------------------------------------

def extract_slide_texts(prs: Presentation) -> list[list[dict]]:
    """
    Returns a list (one per slide) of shape-text records:
        [{"shape_idx": int, "para_idx": int, "text": str}, ...]
    Only non-empty paragraphs are included.
    """
    slides_data = []
    for slide in prs.slides:
        para_records = []
        for s_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    para_records.append({
                        "shape_idx": s_idx,
                        "para_idx": p_idx,
                        "text": text,
                    })
        slides_data.append(para_records)
    return slides_data


# ---------------------------------------------------------------------------
# Translation via Claude (with retry + JSON recovery)
# ---------------------------------------------------------------------------

SYSTEM_PROMPT = """\
You are a professional academic translator specializing in NLP / AI research.
You will receive JSON describing the text content of one PowerPoint slide from a Japanese academic presentation.
Your job is to translate each text entry from Japanese to natural, fluent English suitable for an international conference.

Rules:
- Translate faithfully — do NOT add, remove, or significantly paraphrase content.
- Keep technical terms (model names, dataset names, acronyms) unchanged.
- Keep numbers, symbols, and formulas unchanged.
- Preserve the brevity of slide text — do not expand bullet points into full sentences unless the original is a full sentence.
- Return ONLY a JSON array where each element is {"shape_idx": <int>, "para_idx": <int>, "text": "<translated text>"}.
- Do not include any explanation or markdown fences.
"""

MAX_RETRIES = 3
RETRY_BACKOFF = [2, 5, 15]  # seconds to wait before retry 1, 2, 3


def _strip_fences(raw: str) -> str:
    """Remove accidental markdown code fences from the model response."""
    raw = raw.strip()
    if raw.startswith("```"):
        parts = raw.split("```")
        # parts[1] is the content between first pair of fences
        raw = parts[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return raw.strip()


def _validate_records(translated: list[dict], original: list[dict]) -> list[dict]:
    """
    Cross-check translated records against originals.
    - Must be a list of dicts with required keys.
    - shape_idx / para_idx must match an original record.
    - Falls back to original text for any record that fails validation.
    """
    original_keys = {(r["shape_idx"], r["para_idx"]): r["text"] for r in original}
    validated = []
    for rec in translated:
        if not isinstance(rec, dict):
            continue
        si = rec.get("shape_idx")
        pi = rec.get("para_idx")
        text = rec.get("text", "")
        if si is None or pi is None or not isinstance(text, str):
            continue
        if (si, pi) not in original_keys:
            # Model hallucinated an index — skip
            print(f"    [warn] Translated record has unknown index ({si},{pi}) — skipped",
                  file=sys.stderr)
            continue
        validated.append({"shape_idx": si, "para_idx": pi, "text": text})

    # For any original record missing from the translation, keep original text
    translated_keys = {(r["shape_idx"], r["para_idx"]) for r in validated}
    for rec in original:
        key = (rec["shape_idx"], rec["para_idx"])
        if key not in translated_keys:
            print(f"    [warn] No translation for ({key[0]},{key[1]}) — keeping original",
                  file=sys.stderr)
            validated.append(rec)

    return validated


def translate_slide_batch(
    client: anthropic.Anthropic,
    slide_records: list[dict],
    slide_number: int,
    paper_context: str = "",
) -> list[dict]:
    """
    Translate one slide's text records via Claude.
    Retries on API errors and JSON parse failures.
    Falls back to original text if all retries fail.
    """
    if not slide_records:
        return []

    context_block = (
        f"\n\nPaper context (for terminology reference):\n{paper_context}\n"
        if paper_context else ""
    )
    user_content = (
        f"Slide {slide_number}:{context_block}\n"
        + json.dumps(slide_records, ensure_ascii=False, indent=2)
    )

    last_error = None
    for attempt in range(MAX_RETRIES):
        try:
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2048,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": user_content}],
            )
            raw = _strip_fences(message.content[0].text)
            translated = json.loads(raw)
            if not isinstance(translated, list):
                raise ValueError(f"Expected a JSON array, got {type(translated).__name__}")
            return _validate_records(translated, slide_records)

        except json.JSONDecodeError as e:
            last_error = f"JSON parse error: {e}"
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}",
                  file=sys.stderr)

        except anthropic.RateLimitError as e:
            last_error = f"Rate limit: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except anthropic.APIStatusError as e:
            last_error = f"API error {e.status_code}: {e.message}"
            if e.status_code < 500:
                # 4xx errors (except 429) won't fix themselves — stop retrying
                break
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

        except anthropic.APIConnectionError as e:
            last_error = f"Connection error: {e}"
            wait = RETRY_BACKOFF[min(attempt, len(RETRY_BACKOFF) - 1)]
            print(f"\n    [warn] Attempt {attempt + 1}/{MAX_RETRIES} — {last_error}. "
                  f"Waiting {wait}s...", file=sys.stderr)
            time.sleep(wait)

    print(f"\n    [error] Slide {slide_number}: all retries failed ({last_error}). "
          "Keeping original text.", file=sys.stderr)
    return slide_records  # fall back to originals


# ---------------------------------------------------------------------------
# Write translated text back into the PPTX
# ---------------------------------------------------------------------------

def apply_translations(prs: Presentation, translated_slides: list[list[dict]]) -> None:
    """
    For each translated record, overwrite the paragraph text in-place.
    Strategy: put the full translated string into the first run of the paragraph,
    clear all subsequent runs (preserving formatting objects intact).
    Skips any record whose indices are out of bounds.
    """
    for slide_idx, (slide, records) in enumerate(zip(prs.slides, translated_slides)):
        for rec in records:
            si = rec.get("shape_idx")
            pi = rec.get("para_idx")
            text = rec.get("text", "")

            if si is None or pi is None:
                continue

            # Bounds-check shape index
            if si >= len(slide.shapes):
                print(f"  [warn] Slide {slide_idx + 1}: shape_idx {si} out of range "
                      f"({len(slide.shapes)} shapes) — skipped", file=sys.stderr)
                continue

            shape = slide.shapes[si]
            if not shape.has_text_frame:
                print(f"  [warn] Slide {slide_idx + 1}: shape {si} has no text frame — skipped",
                      file=sys.stderr)
                continue

            paras = shape.text_frame.paragraphs

            # Bounds-check paragraph index
            if pi >= len(paras):
                print(f"  [warn] Slide {slide_idx + 1}: para_idx {pi} out of range "
                      f"({len(paras)} paras in shape {si}) — skipped", file=sys.stderr)
                continue

            para = paras[pi]
            runs = para.runs
            if not runs:
                # No runs (e.g., field placeholder) — skip safely
                continue

            runs[0].text = text
            for run in runs[1:]:
                run.text = ""


# ---------------------------------------------------------------------------
# Checkpoint helpers (resume support)
# ---------------------------------------------------------------------------

def checkpoint_path(output_path: Path) -> Path:
    return output_path.with_suffix(".checkpoint.json")


def load_checkpoint(output_path: Path) -> dict[int, list[dict]] | None:
    cp = checkpoint_path(output_path)
    if cp.exists():
        try:
            data = json.loads(cp.read_text())
            print(f"  Resuming from checkpoint ({len(data)} slides already translated)")
            return {int(k): v for k, v in data.items()}
        except Exception as e:
            print(f"  [warn] Could not read checkpoint: {e} — starting fresh", file=sys.stderr)
    return None


def save_checkpoint(output_path: Path, done: dict[int, list[dict]]) -> None:
    cp = checkpoint_path(output_path)
    cp.write_text(json.dumps(done, ensure_ascii=False, indent=2))


def remove_checkpoint(output_path: Path) -> None:
    cp = checkpoint_path(output_path)
    if cp.exists():
        cp.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Translate Japanese PPTX slides to English.")
    parser.add_argument("input", help="Path to the input .pptx file")
    parser.add_argument("--output", help="Output .pptx path (default: <input>_en.pptx)")
    parser.add_argument("--paper", help="Optional path to the paper PDF for terminology context")
    args = parser.parse_args()

    # Check API key early
    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY is not set.\n"
              "  export ANTHROPIC_API_KEY=sk-ant-...", file=sys.stderr)
        sys.exit(1)

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = (
        Path(args.output) if args.output
        else input_path.with_stem(input_path.stem + "_en")
    )

    paper_context = ""
    if args.paper:
        paper_pdf = Path(args.paper)
        if not paper_pdf.exists():
            print(f"  [warn] Paper PDF not found: {paper_pdf} — skipping context",
                  file=sys.stderr)
        else:
            print(f"Extracting paper context from {paper_pdf.name}...")
            paper_context = extract_pdf_abstract(str(paper_pdf))

    client = anthropic.Anthropic()

    print(f"Loading {input_path.name}...")
    prs = Presentation(str(input_path))
    slides_data = extract_slide_texts(prs)
    total = len(slides_data)

    # Resume from checkpoint if available
    checkpoint = load_checkpoint(output_path) or {}

    print(f"Translating {total} slides...\n")
    translated_slides: list[list[dict]] = []
    try:
        for i, records in enumerate(slides_data, start=1):
            if not records:
                print(f"  Slide {i:02d}/{total}  (empty — skipped)")
                translated_slides.append([])
                continue

            if i in checkpoint:
                print(f"  Slide {i:02d}/{total}  (from checkpoint)")
                translated_slides.append(checkpoint[i])
                continue

            print(f"  Slide {i:02d}/{total}  ({len(records)} text blocks)...",
                  end=" ", flush=True)
            translated = translate_slide_batch(client, records, i, paper_context)
            translated_slides.append(translated)
            checkpoint[i] = translated
            save_checkpoint(output_path, checkpoint)
            print("done")

    except KeyboardInterrupt:
        print("\n\nInterrupted. Progress saved to checkpoint — re-run to resume.", file=sys.stderr)
        sys.exit(1)

    print("\nApplying translations to PPTX...")
    apply_translations(prs, translated_slides)

    prs.save(str(output_path))
    remove_checkpoint(output_path)
    print(f"\nSaved: {output_path}")


if __name__ == "__main__":
    main()
